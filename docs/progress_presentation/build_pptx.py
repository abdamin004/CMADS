"""Build the 3-minute progress presentation (CMADS) as a .pptx file.

Usage:
    python3 build_pptx.py

Output:
    CMADS_Progress_Presentation.pptx  (next to this script)
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SVG = HERE / "agents_block_diagram.svg"
DIAGRAM_PNG = HERE / "agents_block_diagram.png"
OUT_PPTX = HERE / "CMADS_Progress_Presentation.pptx"


# ── Colours ─────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x16, 0xA0, 0x85)
AMBER = RGBColor(0xD3, 0x54, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY_DARK = RGBColor(0x34, 0x49, 0x5E)
GREY_MED = RGBColor(0x7F, 0x8C, 0x8D)
GREY_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)


# ── Helpers ─────────────────────────────────────────────────
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"


def ensure_diagram_png() -> Path:
    """Render the SVG to PNG via headless Chrome.

    qlmanage forces a square aspect ratio; headless Chrome respects the
    SVG's viewBox, so the PNG has the same aspect as the SVG (2.286:1 here).
    """
    if DIAGRAM_PNG.exists():
        return DIAGRAM_PNG
    subprocess.run(
        [
            CHROME, "--headless", "--disable-gpu",
            "--force-device-scale-factor=2",
            f"--screenshot={DIAGRAM_PNG}",
            "--window-size=1600,700",
            "--hide-scrollbars",
            "--default-background-color=FFFFFFFF",
            SVG.as_uri(),
        ],
        check=True, capture_output=True,
    )
    return DIAGRAM_PNG


def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_header_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.name = "Calibri"; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = "Calibri"; r2.font.size = Pt(12); r2.font.color.rgb = GREY_LIGHT
        r2.font.italic = True


def add_textbox(slide, left, top, width, height, text, *,
                size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                italic=False, name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = name; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return tb


def add_bullet_box(slide, left, top, width, height, bullets,
                   size=14, color=GREY_DARK, gap=Pt(6)):
    """bullets: list of (text, bold?) OR list of str."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = gap
        r = p.add_run(); r.text = "•  " + text
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_footer(slide, page: str):
    add_textbox(slide, Inches(0.3), Inches(7.1), Inches(8), Inches(0.3),
                "CMADS · Bachelor Thesis · Abdelrahman · April 2026",
                size=10, color=GREY_MED, italic=True)
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
                page, size=10, color=GREY_MED, align=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=NAVY, card_color=WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = card_color
    card.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0); card.line.width = Pt(0.75)
    card.shadow.inherit = False
    # Value
    tb = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.name = "Calibri"; r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = value_color
    # Label
    tb2 = slide.shapes.add_textbox(left, top + Inches(0.85), width, Inches(0.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.name = "Calibri"; r2.font.size = Pt(11); r2.font.color.rgb = GREY_DARK


def add_speaker_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ── Build deck ──────────────────────────────────────────────
def build():
    diagram = ensure_diagram_png()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1 · Title ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, PAPER)

    # Left colour band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.6), Inches(11.5), Inches(0.5),
                "BACHELOR THESIS · PROGRESS UPDATE", size=14, bold=True, color=AMBER)

    add_textbox(s, Inches(1), Inches(2.1), Inches(11.5), Inches(1.3),
                "CMADS", size=64, bold=True, color=NAVY)

    add_textbox(s, Inches(1), Inches(3.3), Inches(11.5), Inches(0.9),
                "Clinical Multi-Agent Decisioning System",
                size=30, bold=True, color=BLUE)

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.5), Inches(0.6),
                "A LangGraph pipeline for differential diagnosis and NICE-guideline "
                "treatment planning on synthetic patients.",
                size=16, italic=True, color=GREY_DARK)

    # Divider
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.1),
                              Inches(3), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = AMBER; line.line.fill.background()

    add_textbox(s, Inches(1), Inches(5.3), Inches(11.5), Inches(0.4),
                "Presented by  Abdelrahman", size=16, bold=True, color=GREY_DARK)
    add_textbox(s, Inches(1), Inches(5.75), Inches(11.5), Inches(0.4),
                "Supervisor meeting · April 2026", size=14, color=GREY_MED)

    add_textbox(s, Inches(1), Inches(6.6), Inches(11.5), Inches(0.3),
                "Agenda   ·   Methodology   ·   Results on 270 patients   ·   "
                "Model comparison   ·   Multi-level memory (NEW)   ·   Next steps",
                size=11, color=GREY_MED, italic=True)

    add_speaker_notes(s, """
Good afternoon. Today I'll briefly cover three things:
one, the methodology I settled on; two, the results on 270 verified patients;
three, a head-to-head comparison between a large general model and a
medical-fine-tuned model that I finished last week. Target: three minutes.
The research question is whether a coordinated group of specialised LLM
agents can replicate clinical reasoning given only structured EHR and lab data.
""")

    # ── Slide 2 · Multi-Agent Pipeline ─────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(s, "Multi-Agent Pipeline",
                   "7 specialised agents · 6 stages · LangGraph StateGraph · "
                   "shared TypedDict state")

    # Diagram placeholder — a dashed frame the user drops the image into.
    # Fit by height → height 5.55 in, width 5.55 × 2.286 = 12.69 in.
    diagram_h = 5.55
    diagram_w = diagram_h * (1600 / 700)   # ≈ 12.69 in
    slide_w = 13.333
    diagram_left = (slide_w - diagram_w) / 2  # centre horizontally
    diagram_top = 1.05

    placeholder = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(diagram_left), Inches(diagram_top),
        Inches(diagram_w), Inches(diagram_h),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF7)
    placeholder.line.color.rgb = RGBColor(0x99, 0xA3, 0xA4)
    placeholder.line.width = Pt(1.5)
    placeholder.line.dash_style = 7  # dash

    add_textbox(
        s,
        Inches(diagram_left), Inches(diagram_top + diagram_h/2 - 0.6),
        Inches(diagram_w), Inches(0.5),
        "Insert diagram here",
        size=22, bold=True, color=GREY_MED, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(diagram_left), Inches(diagram_top + diagram_h/2 - 0.1),
        Inches(diagram_w), Inches(0.4),
        "agents_block_diagram.svg   ·   .pdf   ·   _hires.png (6400 × 2800)",
        size=13, italic=True, color=GREY_MED, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(diagram_left), Inches(diagram_top + diagram_h/2 + 0.4),
        Inches(diagram_w), Inches(0.4),
        "Recommended placement: " + f"{diagram_w:.2f} in × {diagram_h:.2f} in, "
        f"top-left at ({diagram_left:.2f}, {diagram_top:.2f})",
        size=11, italic=True, color=GREY_MED, align=PP_ALIGN.CENTER,
    )

    # Thin caption strip under the placeholder
    cap_top = diagram_top + diagram_h + 0.1
    cap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.4), Inches(cap_top),
                             Inches(12.55), Inches(0.3))
    cap.fill.solid(); cap.fill.fore_color.rgb = NAVY; cap.line.fill.background()
    add_textbox(s, Inches(0.5), Inches(cap_top + 0.03),
                Inches(12.35), Inches(0.28),
                "Fan-out at Stage 1   ·   adaptive self-critique at Stage 2   ·   "
                "adversarial review at Stage 3   ·   "
                "treatment only on DIRECT matches at Stage 6",
                size=10.5, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    add_footer(s, "2 / 6")
    add_speaker_notes(s, """
This is the thesis contribution: seven agents orchestrated with LangGraph,
reading and writing a single shared TypedDict state.
Stage one — EHR Analyst and Lab Interpreter — fans out in parallel from the
start node; merge reducers on the state channel prevent them from overwriting
each other.
Stage two, Diagnostic Reasoning, runs an adaptive self-critique loop: it
generates a differential, critiques itself, and refines up to three rounds,
stopping early when confidence crosses seventy-five.
Stage three, the Clinical Reviewer, provides an independent second opinion —
per-diagnosis CONFIRM, MODIFY, or REJECT verdicts.
Stage four, the Refiner, merges both perspectives into the final differential.
Stage five is LLM-as-judge evaluation against Synthea ground truth.
Stage six, Treatment Planning, fires only on DIRECT matches and retrieves
NICE guideline passages from Qdrant.
Every agent follows the same five-component blueprint shown at the bottom —
input gate, prompt, LLM, parse, output gate — which is what makes the pipeline
config-driven and any single agent independently testable.
""")

    # ── Slide 3 · Results on 160 patients ─────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(s, "Results — 160-patient cohort",
                   "Actual MAS pipeline runs · GPT-OSS 120B via Groq · "
                   "Qwen3 32B judge · 8 disease categories")

    # KPI row (four headline numbers from aggregate_160.json)
    kpi_top = Inches(1.15); kpi_w = Inches(2.85); kpi_h = Inches(1.35)
    add_kpi_card(s, Inches(0.4),  kpi_top, kpi_w, kpi_h, "74 %",
                 "DIRECT match (118 / 160)", value_color=GREEN)
    add_kpi_card(s, Inches(3.45), kpi_top, kpi_w, kpi_h, "88 %",
                 "Found rate — DIRECT + INDIRECT", value_color=NAVY)
    add_kpi_card(s, Inches(6.50), kpi_top, kpi_w, kpi_h, "60 %",
                 "Rank-1 placement when found", value_color=BLUE)
    add_kpi_card(s, Inches(9.55), kpi_top, kpi_w, kpi_h, "~193 s",
                 "Avg pipeline time / patient", value_color=AMBER)

    # Per-disease table for 160 patients
    table_left = Inches(0.4); table_top = Inches(2.8)
    table_w = Inches(7.9); table_h = Inches(4.0)

    rows = [
        ("Disease", "n", "DIRECT", "INDIRECT", "MISS", "Found"),
        ("End-stage renal disease",      "51", "41 · 80 %", "9",  "1",  "98 %"),
        ("Metabolic syndrome X",         "32", "29 · 91 %", "0",  "3",  "91 %"),
        ("Essential hypertension",       "25", "16 · 64 %", "2",  "7",  "72 %"),
        ("Ischemic heart disease",       "22", "16 · 73 %", "3",  "3",  "86 %"),
        ("CKD stage 3",                  "14",  "8 · 57 %", "5",  "1",  "93 %"),
        ("Diabetes mellitus type 2",      "8",  "4 · 50 %", "1",  "3",  "63 %"),
        ("Chronic heart failure",         "5",  "3 · 60 %", "0",  "2",  "60 %"),
        ("CKD stage 2",                   "3",  "1 · 33 %", "2",  "0", "100 %"),
        ("TOTAL",                       "160", "118 · 74 %", "22", "20", "88 %"),
    ]
    table = s.shapes.add_table(len(rows), 6,
                               table_left, table_top, table_w, table_h).table
    table.columns[0].width = Inches(2.5)
    for i in range(1, 6):
        table.columns[i].width = Inches(1.08)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(11) if i not in (0, len(rows) - 1) else Pt(12)
                    r.font.bold = (i == 0 or i == len(rows) - 1)
                    if i == 0:
                        r.font.color.rgb = WHITE
                    elif i == len(rows) - 1:
                        r.font.color.rgb = WHITE
                    else:
                        r.font.color.rgb = GREY_DARK
                        if j == 2:  # DIRECT column coloured by performance
                            pct_text = val.split("·")[-1].strip().rstrip("%")
                            try:
                                pct = float(pct_text)
                                if pct >= 70:
                                    r.font.color.rgb = GREEN
                                    r.font.bold = True
                                elif pct >= 50:
                                    r.font.color.rgb = AMBER
                                    r.font.bold = True
                                else:
                                    r.font.color.rgb = RED
                            except ValueError:
                                pass
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif i == len(rows) - 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = GREY_DARK
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF7)

    # Right-side findings panel
    fpanel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(8.5), Inches(2.8), Inches(4.6), Inches(4.0))
    fpanel.fill.solid(); fpanel.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
    fpanel.line.color.rgb = RGBColor(0xD4, 0xAC, 0x0D); fpanel.line.width = Pt(0.75)
    add_textbox(s, Inches(8.7), Inches(2.9), Inches(4.3), Inches(0.4),
                "Key findings", size=14, bold=True, color=RGBColor(0x7D, 0x66, 0x08))
    add_bullet_box(s, Inches(8.7), Inches(3.35), Inches(4.3), Inches(3.4), [
        ("ESRD leads at 98 % found · 80 % DIRECT (n=51).", True),
        ("Metabolic Syndrome solved — 91 % DIRECT (n=32) once composite-dx prompts were tuned.", True),
        ("IHD jumped from 0 % → 73 % DIRECT after lab-vs-EHR rebalancing.", True),
        ("Hypertension dipped to 64 % DIRECT — more nuanced cases in the expanded cohort.", False),
        ("Smallest cohorts (CHF, CKD-2) have widest confidence intervals — next to expand.", False),
    ], size=11, color=GREY_DARK, gap=Pt(5))

    add_footer(s, "3 / 6")
    add_speaker_notes(s, """
These are the actual numbers across all one hundred sixty patients that have
been run through the pipeline so far. Seventy-four percent DIRECT match —
118 out of 160 — and an eighty-eight percent found rate if we include INDIRECT
matches. When the system finds the right disease, sixty percent of the time
it lands at rank one.
The per-disease picture shows the improvements from the last two iterations:
metabolic syndrome, which was zero on the early fifty-patient run, is now at
ninety-one percent DIRECT after composite-diagnosis prompt tuning. Ischemic
heart disease went from zero to seventy-three percent after rebalancing how
the Diagnostic agent weighs EHR history against dominant renal labs.
ESRD at eighty percent DIRECT with fifty-one patients is the strongest signal.
The categories with smallest n — CHF and CKD stage 2 — are the ones to expand
before the thesis writeup.
""")

    # ── Slide 4 · Model comparison ────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(s, "Model comparison — GPT-OSS 120B vs Med42 70B",
                   "20 patients · Batch 4 (seed 42) · same pipeline · same Qwen3 32B judge")

    # Comparison table
    rows = [
        ("Metric", "GPT-OSS 120B (Groq)", "Med42 70B (Ollama, local)"),
        ("DIRECT match", "50 %", "25 %"),
        ("Found rate (D + I)", "80 %", "75 %"),
        ("Head-to-head wins", "8", "2   (10 ties)"),
        ("Avg time per patient", "~2 min", "~27 min  (13× slower)"),
        ("Cost per patient", "$0.06", "$0.00"),
        ("Agent success rate", "100 %", "40 – 85 %   (JSON / timeout)"),
    ]
    tbl_left = Inches(0.4); tbl_top = Inches(1.1)
    tbl_w = Inches(7.3); tbl_h = Inches(3.2)
    table = s.shapes.add_table(len(rows), 3, tbl_left, tbl_top, tbl_w, tbl_h).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(2.45)
    table.columns[2].width = Inches(2.45)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(12)
                    r.font.bold = (i == 0)
                    if i == 0:
                        r.font.color.rgb = WHITE
                    elif j == 1:   # GPT-OSS column
                        r.font.color.rgb = NAVY
                        r.font.bold = True
                    elif j == 2:   # Med42
                        r.font.color.rgb = GREY_DARK
                    else:
                        r.font.color.rgb = GREY_DARK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF7)

    # Findings panel
    fpanel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(7.9), Inches(1.1), Inches(5.2), Inches(5.4))
    fpanel.fill.solid(); fpanel.fill.fore_color.rgb = WHITE
    fpanel.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0); fpanel.line.width = Pt(0.75)
    add_textbox(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(0.45),
                "Four findings", size=15, bold=True, color=AMBER)

    findings = [
        ("1 · General beats medical fine-tune at this size.", True),
        ("GPT-OSS dominates composite diagnoses — Metabolic Syndrome 5/6 vs 1/6.", False),
        ("2 · Med42 edges ahead on plain hypertension.", True),
        ("Medical pre-training helps vitals-based dx (2/4 DIRECT vs 1/4).", False),
        ("3 · Self-evaluation is inflated.", True),
        ("Med42 scored itself 45 % DIRECT; Qwen3 judge cut it to 25 %.", False),
        ("→ A fixed third-party evaluator is methodologically essential.", False),
        ("4 · Neither model resolves CKD staging.", True),
        ("0 % DIRECT on stage 2/3 for both — structural gap, not a model gap.", False),
    ]
    add_bullet_box(s, Inches(8.15), Inches(1.75), Inches(4.95), Inches(4.7),
                   findings, size=11.5, color=GREY_DARK, gap=Pt(5))

    # Extra row under table
    add_textbox(s, Inches(0.4), Inches(4.5), Inches(7.3), Inches(0.4),
                "Head-to-head: GPT-OSS wins 8, Med42 wins 2, tied on 10.",
                size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(0.4), Inches(4.95), Inches(7.3), Inches(1.5),
                "Per-disease highlights   ·   Metabolic Syndrome: GPT-OSS 83 % vs Med42 17 %   "
                "·   Hypertension: GPT-OSS 25 % vs Med42 50 %   ·   IHD: GPT-OSS 75 % vs Med42 25 %",
                size=11.5, italic=True, color=GREY_DARK)

    add_footer(s, "4 / 6")
    add_speaker_notes(s, """
This is the experiment I finished most recently. I ran the same 20 patients
through the same pipeline twice, changing only the LLM, and evaluated both
with Qwen3 32B as a neutral judge.
GPT-OSS wins 50 to 25 percent on DIRECT matches, and eight to two head-to-head.
The most interesting finding is point three: when Med42 evaluated its own
output it scored itself at 45 percent, but under the neutral judge it dropped
to 25 percent. That is a methodological result worth keeping — any
self-evaluation number in this domain is untrustworthy.
Med42's medical fine-tuning does help on simple vitals-based conditions like
hypertension, and it costs zero dollars, but at thirteen times the wall-clock.
""")

    # ── Slide 5 · Multi-Level Memory (NEW · supervisor request) ──────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s,
        "Multi-Level Memory  ·  supervisor request",
        "Four-tier memory layered on the existing LangGraph state  ·  "
        "MEMORY_ENABLED flag toggles for clean A/B"
    )

    # Left half: four tier cards stacked vertically.
    tier_left = Inches(0.4)
    tier_w = Inches(7.4)
    tier_h = Inches(1.18)
    tier_gap = Inches(0.18)
    tier_top0 = Inches(1.15)

    tiers = [
        ("Tier 1 · WORKING",
         "per-agent · invocation-local",
         "Confidence trajectory + critique trail inside the "
         "Diagnostic adaptive loop.",
         AMBER),
        ("Tier 2 · EPISODIC",
         "current run · in-state, persisted",
         "Typed timeline of events. Reviewer & Refiner read the "
         "reasoning chain, not just final JSON.",
         BLUE),
        ("Tier 3 · SEMANTIC",
         "cross-session · on disk",
         "Per-disease aggregate stats (DIRECT/INDIRECT/MISS, "
         "rank-1, evidence patterns) written by a new Stage 7.",
         TEAL),
        ("Tier 4 · PROCEDURAL",
         "long-term · static (existing Qdrant)",
         "NICE guidelines, now wrapped in a uniform recall(query) "
         "API so all four tiers feel the same.",
         NAVY),
    ]

    for i, (title, scope, body, accent) in enumerate(tiers):
        top = tier_top0 + (tier_h + tier_gap) * i
        # Left accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 tier_left, top, Inches(0.12), tier_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        # Card
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  tier_left + Inches(0.12), top,
                                  tier_w - Inches(0.12), tier_h)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
        card.line.width = Pt(0.5)
        # Title
        add_textbox(s,
                    tier_left + Inches(0.3), top + Inches(0.08),
                    Inches(4.5), Inches(0.4),
                    title, size=15, bold=True, color=accent)
        # Scope (right-aligned on the same row)
        add_textbox(s,
                    tier_left + Inches(4.8), top + Inches(0.13),
                    Inches(2.5), Inches(0.4),
                    scope, size=11, italic=True,
                    color=GREY_MED, align=PP_ALIGN.RIGHT)
        # Body
        add_textbox(s,
                    tier_left + Inches(0.3), top + Inches(0.5),
                    Inches(7.0), Inches(0.65),
                    body, size=12.5, color=GREY_DARK)

    # Right half: A/B comparison + impact panel.
    panel_left = Inches(8.05)
    panel_top = Inches(1.15)
    panel_w = Inches(5.05)
    panel_h = Inches(5.65)

    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        panel_left, panel_top, panel_w, panel_h,
    )
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
    panel.line.color.rgb = RGBColor(0xD4, 0xAC, 0x0D); panel.line.width = Pt(0.75)

    add_textbox(s, panel_left + Inches(0.2), panel_top + Inches(0.15),
                panel_w - Inches(0.4), Inches(0.4),
                "Before / after on Batch 4 (n=20)",
                size=14, bold=True,
                color=RGBColor(0x7D, 0x66, 0x08))
    add_textbox(s, panel_left + Inches(0.2), panel_top + Inches(0.55),
                panel_w - Inches(0.4), Inches(0.4),
                "GPT-OSS 120B · same pipeline · same Qwen3 judge",
                size=10.5, italic=True, color=GREY_MED)

    ab_rows = [
        ("Metric", "OFF", "ON"),
        ("DIRECT match", "—", "—"),
        ("Found rate (D + I)", "—", "—"),
        ("Avg time / patient", "—", "—"),
        ("Rank-1 when found", "—", "—"),
    ]
    ab_top = panel_top + Inches(1.05)
    ab_left = panel_left + Inches(0.2)
    ab_w = panel_w - Inches(0.4)
    ab_h = Inches(2.1)
    table = s.shapes.add_table(len(ab_rows), 3,
                               ab_left, ab_top, ab_w, ab_h).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(1.25)
    table.columns[2].width = Inches(1.25)
    for i, row in enumerate(ab_rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Calibri"; r.font.size = Pt(11.5)
                    r.font.bold = (i == 0)
                    if i == 0:
                        r.font.color.rgb = WHITE
                    else:
                        r.font.color.rgb = GREY_DARK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF7)

    add_textbox(s, panel_left + Inches(0.2), ab_top + ab_h + Inches(0.05),
                panel_w - Inches(0.4), Inches(0.45),
                "Numbers populated after the two 20-patient A/B runs finish.",
                size=10, italic=True, color=GREY_MED)

    add_textbox(s, panel_left + Inches(0.2),
                panel_top + Inches(3.85),
                panel_w - Inches(0.4), Inches(0.4),
                "What it changes",
                size=13, bold=True,
                color=RGBColor(0x7D, 0x66, 0x08))
    add_bullet_box(
        s,
        panel_left + Inches(0.2), panel_top + Inches(4.3),
        panel_w - Inches(0.4), Inches(1.3),
        [
            ("Reviewer cites specific critique rounds from "
             "Diagnostic, not just final dx.", False),
            ("Refiner reads cross-session priors for "
             "candidate diseases.", False),
            ("Each run consolidates outcomes into Tier-3 — "
             "next run starts smarter.", False),
            ("All four tiers covered by 24 unit tests.", True),
        ],
        size=10.5, color=GREY_DARK, gap=Pt(3),
    )

    # Caption strip under the tier cards.
    cap_top = tier_top0 + (tier_h + tier_gap) * 4 + Inches(0.05)
    cap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             tier_left, cap_top, tier_w, Inches(0.45))
    cap.fill.solid(); cap.fill.fore_color.rgb = NAVY
    cap.line.fill.background()
    add_textbox(s, tier_left + Inches(0.2), cap_top + Inches(0.06),
                tier_w - Inches(0.4), Inches(0.35),
                "All four tiers reached through one MemoryManager facade  ·  "
                "writes flow through LangGraph reducers  ·  no LLM coupling",
                size=11, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    add_footer(s, "5 / 6")
    add_speaker_notes(s, """
Last week you asked for a multi-level memory system so the agents could share
session context, not only their final JSON outputs. This is the implementation.
Four tiers — working, episodic, semantic, procedural — inspired by the CoALA
cognitive-architecture paper but specialised to this clinical workflow.
Working memory is per-agent scratch, used by the Diagnostic agent to track its
own confidence trajectory across critique rounds. Episodic is a typed timeline
of decisions and confidence checks the Reviewer and Refiner can read so they
reason about the path to the diagnosis, not only the diagnosis. Semantic is a
small JSON file on disk that accumulates per-disease statistics across runs;
a new Stage seven, Memory Consolidation, writes into it after every patient.
Procedural is the existing NICE guidelines via Qdrant, now exposed through the
same uniform API.
The whole subsystem is gated by a single MEMORY_ENABLED flag, which is how
the before-and-after run on Batch 4 stays clean. The two 20-patient runs are
queued; numbers go in the empty cells before the meeting.
""")

    # ── Slide 6 · Status summary ──────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(s, "Status summary",
                   "Where the project stands today")

    # Central hero card
    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(1.25),
                              Inches(12.1), Inches(4.9))
    hero.fill.solid(); hero.fill.fore_color.rgb = NAVY; hero.line.fill.background()

    # Three headline numbers as big tiles across the top of the hero card
    kpi_w = Inches(3.7); kpi_h = Inches(1.6)
    add_kpi_card(s, Inches(0.95),  Inches(1.55), kpi_w, kpi_h,
                 "160", "Patients run end-to-end", value_color=NAVY)
    add_kpi_card(s, Inches(4.85),  Inches(1.55), kpi_w, kpi_h,
                 "74 %", "DIRECT diagnosis match", value_color=GREEN)
    add_kpi_card(s, Inches(8.75),  Inches(1.55), kpi_w, kpi_h,
                 "88 %", "Found rate (DIRECT + INDIRECT)", value_color=AMBER)

    # Status bullets on the dark hero panel
    add_textbox(s, Inches(1.1), Inches(3.4), Inches(11.1), Inches(0.45),
                "WHERE WE STAND", size=14, bold=True, color=AMBER)

    bullet_items = [
        "7-agent LangGraph pipeline is stable — parallel Stage 1, adaptive Stage 2, adversarial Stage 3.",
        "Provider-agnostic LLM adapter — switching Groq ↔ Ollama ↔ OpenAI is one .env line.",
        "End-to-end evaluation against Synthea ground truth is reproducible and fully automated.",
        "Model comparison protocol validated on GPT-OSS vs Med42 (50 % vs 25 % DIRECT).",
        "160 / 270 verified patients processed · ~193 s / patient · ≈ $0.06 / patient on Groq.",
    ]
    tb = s.shapes.add_textbox(Inches(1.1), Inches(3.9),
                              Inches(11.1), Inches(2.1))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(bullet_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "✓  " + item
        r.font.name = "Calibri"; r.font.size = Pt(13)
        r.font.color.rgb = WHITE

    # Thank-you strip at the bottom
    add_textbox(s, Inches(0.4), Inches(6.4), Inches(12.6), Inches(0.5),
                "Thank you — questions?", size=22, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(s, "6 / 6")
    add_speaker_notes(s, """
To summarise. The seven-agent LangGraph pipeline is stable and reproducible.
One hundred sixty of the two hundred seventy verified patients have been run
end-to-end with seventy-four percent DIRECT match and eighty-eight percent
found rate. The model-comparison protocol is validated, and the multi-level
memory feature you requested is wired in, tested, and ready to evaluate.
Thank you — happy to take questions.
""")

    # Save
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    build()
