"""Build the CMADS Results & Comparison deck.

9 slides. Numbers come from compute_metrics.compute(). No imports
from docs/final_presentation/. Run:

    python3 docs/results_presentation/build_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import compute_metrics

HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
TH_IMG = REPO / "thesis" / "images"
OUT = HERE / "CMADS_Results_Presentation.pptx"

# ── Palette ─────────────────────────────────────────────────────────
NAVY = RGBColor(0x14, 0x2B, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x16, 0xA0, 0x85)
AMBER = RGBColor(0xD3, 0x54, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY_DARK = RGBColor(0x34, 0x49, 0x5E)
GREY_MED = RGBColor(0x7F, 0x8C, 0x8D)
GREY_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── Layout primitives ───────────────────────────────────────────────

def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background(); bg.shadow.inherit = False


def add_header(slide, title: str, subtitle: str | None = None, accent=NAVY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.7))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = GREY_LIGHT


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False,
                color=GREY_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def add_metric_tile(slide, left, top, width, height, value: str, label: str,
                    value_color=NAVY, bg=WHITE):
    """A big-number tile: value on top, small label under it."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = GREY_LIGHT; box.line.width = Pt(1)
    box.shadow.inherit = False
    tb = box.text_frame; tb.word_wrap = True
    tb.margin_left = tb.margin_right = Inches(0.05)
    p1 = tb.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = value
    r1.font.size = Pt(34); r1.font.bold = True; r1.font.color.rgb = value_color
    p2 = tb.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11); r2.font.color.rgb = GREY_MED


def add_table(slide, left, top, width, height, headers, rows,
              *, header_bg=NAVY, header_fg=WHITE, body_fg=GREY_DARK,
              first_col_bold=True):
    """Plain table with a coloured header row."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(height))
    table = tbl_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame; tf.text = ""
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = header_fg
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else PAPER
            tf = cell.text_frame; tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(11); r.font.color.rgb = body_fg
            r.font.bold = first_col_bold and j == 0
    return table


def add_image(slide, path: Path, left, top, *, width=None, height=None):
    kwargs = {}
    if width: kwargs["width"] = Inches(width)
    if height: kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


# ── Slide functions (filled in by later tasks) ──────────────────────

def slide_title(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.5))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    add_textbox(slide, 0.6, 2.0, 12.0, 1.2,
                "CMADS — Multi-Agent Systems for AI Clinical Decisioning",
                size=36, bold=True, color=NAVY)
    add_textbox(slide, 0.6, 3.2, 12.0, 0.7,
                "via Automation Workflows",
                size=28, color=GREY_DARK)
    add_textbox(slide, 0.6, 4.3, 12.0, 0.6,
                "100-patient multi-level memory results  ·  doctor console  ·  literature comparison",
                size=16, color=TEAL)
    add_textbox(slide, 0.6, 6.4, 12.0, 0.5,
                "Abdelrahman Mohamed Amin  ·  Supervisor: Dr. Shereen Moataz Afifi  ·  GUC  ·  May 2026",
                size=13, color=GREY_MED)


def slide_what_is_cmads(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "What CMADS is",
               "7-agent LangGraph pipeline + 4-tier memory subsystem")

    sysarch = TH_IMG / "ch3_system_architecture.png"
    if sysarch.exists():
        add_image(slide, sysarch, 0.3, 1.2, width=8.4)
    else:
        add_textbox(slide, 0.5, 1.5, 8.0, 1.0,
                    "[diagram missing: ch3_system_architecture.png]",
                    size=12, color=RED)

    bullets = [
        ("Synthetic patient generation",
         "Synthea → Bronze → Silver → Gold medallion pipeline; one cohort, no PHI."),
        ("7-agent LangGraph pipeline",
         "EHR + Labs → diagnosis → review → refine → evaluate → NICE treatment."),
        ("Vector database (Qdrant)",
         "BioLORD-2023 embeddings of past cases + NICE guidelines; RAG at recall time."),
        ("4-tier memory subsystem",
         "Working · Episodic · Semantic · Case-based — written every patient."),
        ("Open-source LLMs",
         "GPT-OSS-120B reasoning + Qwen3-32B judge via Groq — reproducible under $30/1k."),
        ("Doctor-facing console",
         "Agent inspector, similar-case browser, treatment-safety review, persisted verdicts."),
    ]
    left = 8.9
    top0 = 1.25
    line_h = 0.9
    for i, (head, body) in enumerate(bullets):
        top = top0 + i * line_h
        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(left), Inches(top + 0.12),
                                      Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        add_textbox(slide, left + 0.28, top - 0.02, 4.0, 0.4,
                    head, size=12, bold=True, color=NAVY)
        add_textbox(slide, left + 0.28, top + 0.32, 4.0, 0.6,
                    body, size=10, color=GREY_DARK)

    add_textbox(slide, 0.4, 6.85, 12.5, 0.4,
                "Stage 1 (parallel) → Diagnostic loop → Reviewer → Refiner → Evaluator → Treatment (DIRECT only) → Memory consolidation",
                size=10, color=GREY_MED, align=PP_ALIGN.CENTER)


def slide_headline(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Results — multi-level memory vs single-level baseline",
               "Memory ON (n=100, batch_3 + batch_4) vs memory OFF (n=160, mas_results/)")

    mem = metrics["cohorts"]["combined_100"]
    base = metrics["cohorts"]["single_level_baseline"]

    # Two side-by-side metric cards
    cards = [
        ("Single-level baseline",
         "Memory OFF · 160 patients",
         base, BLUE),
        ("Multi-level memory",
         "Memory ON · 100 patients",
         mem, AMBER),
    ]
    for i, (title, sub, agg, accent) in enumerate(cards):
        left = 0.6 + i * 6.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(1.2),
                                      Inches(5.9), Inches(3.6))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = accent; box.line.width = Pt(1.5)
        add_textbox(slide, left + 0.25, 1.35, 5.5, 0.4, title,
                    size=16, bold=True, color=accent)
        add_textbox(slide, left + 0.25, 1.78, 5.5, 0.4, sub,
                    size=11, color=GREY_MED)
        lines = [
            ("DIRECT",          f"{agg['DIRECT_pct']:.1f}%   ({agg['DIRECT']}/{agg['n']})"),
            ("Found (D + I)",   f"{agg['found_pct']:.1f}%   ({agg['found']}/{agg['n']})"),
            ("MISS",            f"{agg['MISS_pct']:.1f}%   ({agg['MISS']}/{agg['n']})"),
            ("Rank-1 in found", f"{agg['rank1_in_found_pct']:.1f}%   ({agg['rank1_in_found']}/{agg['found']})"),
        ]
        for j, (k, v) in enumerate(lines):
            top = 2.3 + j * 0.6
            add_textbox(slide, left + 0.3, top, 2.7, 0.45, k,
                        size=12, bold=True, color=GREY_DARK)
            add_textbox(slide, left + 3.0, top, 2.8, 0.45, v,
                        size=14, bold=True, color=accent)

    # Delta strip
    d_dir = mem["DIRECT_pct"] - base["DIRECT_pct"]
    d_found = mem["found_pct"] - base["found_pct"]
    d_miss = mem["MISS_pct"] - base["MISS_pct"]

    add_textbox(slide, 0.6, 5.0, 12.1, 0.5,
                "Δ (Multi-level − Single-level)",
                size=14, bold=True, color=GREY_DARK)
    deltas = [
        ("DIRECT", f"{d_dir:+.1f} pp",  RED if d_dir < 0 else GREEN),
        ("Found",  f"{d_found:+.1f} pp", GREEN if d_found > 0 else RED),
        ("MISS",   f"{d_miss:+.1f} pp",  GREEN if d_miss < 0 else RED),
    ]
    for i, (k, v, col) in enumerate(deltas):
        left = 0.6 + i * 4.1
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(left), Inches(5.5),
                                       Inches(3.9), Inches(0.7))
        chip.fill.solid(); chip.fill.fore_color.rgb = WHITE
        chip.line.color.rgb = GREY_LIGHT
        add_textbox(slide, left + 0.2, 5.6, 1.5, 0.5, k,
                    size=14, bold=True, color=GREY_DARK)
        add_textbox(slide, left + 1.7, 5.6, 2.1, 0.5, v,
                    size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)

    add_textbox(slide, 0.6, 6.5, 12.1, 0.9,
                "Memory broadens recall (Found +7.5 pp, MISS −6.7 pp) but "
                "shifts confirmed matches into the related-but-not-exact bucket "
                "(DIRECT −24.1 pp). Different patient subsets per arm — "
                "comparison is descriptive, not a paired test.",
                size=11, color=GREY_MED)


def slide_regime_split(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Split by regime",
               "Cold-start (no same-cohort priors) vs warmed (50 batch_4 priors)")
    cold = metrics["cohorts"]["batch_3_cold_start"]
    warm = metrics["cohorts"]["batch_4_warmed"]
    cards = [
        ("COLD-START — batch_3 (50)",
         "0 same-cohort priors at run start.",
         cold, BLUE),
        ("WARMED — batch_4 (50)",
         "Qdrant already held 50 batch_4 patients from a prior run.",
         warm, AMBER),
    ]
    for i, (title, sub, agg, accent) in enumerate(cards):
        left = 0.7 + i * 6.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(1.3),
                                      Inches(5.8), Inches(4.3))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = accent; box.line.width = Pt(1.5)
        add_textbox(slide, left + 0.2, 1.45, 5.5, 0.5, title,
                    size=16, bold=True, color=accent)
        add_textbox(slide, left + 0.2, 1.95, 5.5, 0.5, sub,
                    size=11, color=GREY_MED)
        lines = [
            ("DIRECT",          f"{agg['DIRECT_pct']:.0f}%   ({agg['DIRECT']}/{agg['n']})"),
            ("Found",           f"{agg['found_pct']:.0f}%   ({agg['found']}/{agg['n']})"),
            ("Rank-1 in found", f"{agg['rank1_in_found_pct']:.0f}%   ({agg['rank1_in_found']}/{agg['found']})"),
        ]
        for j, (k, v) in enumerate(lines):
            top = 2.7 + j * 0.85
            add_textbox(slide, left + 0.3, top, 2.4, 0.5, k,
                        size=13, bold=True, color=GREY_DARK)
            add_textbox(slide, left + 2.8, top, 2.9, 0.5, v,
                        size=18, bold=True, color=accent)
    leak = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.7), Inches(5.85),
                                   Inches(11.9), Inches(0.95))
    leak.fill.solid(); leak.fill.fore_color.rgb = GREY_LIGHT
    leak.line.fill.background()
    add_textbox(slide, 0.9, 5.95, 11.5, 0.45,
                "Cohort-leakage estimate (notes/experiments.md, 2026-05-11):",
                size=12, bold=True, color=GREY_DARK)
    add_textbox(slide, 0.9, 6.35, 11.5, 0.45,
                "~6 pp of the batch_4 DIRECT gain and ~6 pp of the Found gain "
                "come from same-cohort priors, not algorithm.",
                size=12, color=GREY_DARK)


def slide_paired_mcnemar(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Paired McNemar A/B (n = 20)",
               "Same 20 patients, memory toggled · the only controlled test in the project")
    p = metrics["paired_mcnemar"]
    c = p["contingency"]
    headers = ["", "ON · DIRECT", "ON · not-DIRECT"]
    rows = [
        ["OFF · DIRECT",     c["both_DIRECT"],    c["only_OFF_DIRECT"]],
        ["OFF · not-DIRECT", c["only_ON_DIRECT"], c["neither_DIRECT"]],
    ]
    add_table(slide, 0.7, 1.6, 6.5, 2.9, headers, rows)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(7.7), Inches(1.6),
                                    Inches(5.0), Inches(2.9))
    panel.fill.solid(); panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = GREY_LIGHT
    add_textbox(slide, 7.9, 1.7, 4.6, 0.5,
                "DIRECT rate", size=13, bold=True, color=GREY_DARK)
    add_textbox(slide, 7.9, 2.2, 4.6, 0.7,
                f"{p['off_direct_rate']*100:.0f}%  →  {p['on_direct_rate']*100:.0f}%",
                size=28, bold=True, color=NAVY)
    add_textbox(slide, 7.9, 3.0, 4.6, 0.5,
                "(8/20 → 9/20 · +5 pp point estimate)",
                size=11, color=GREY_MED)
    add_textbox(slide, 7.9, 3.6, 4.6, 0.5,
                "Exact McNemar (two-sided)",
                size=13, bold=True, color=GREY_DARK)
    add_textbox(slide, 7.9, 4.0, 4.6, 0.5,
                f"p = {p['mcnemar_p_two_sided']:.1f}   ·   discordant pairs: "
                f"{c['only_OFF_DIRECT'] + c['only_ON_DIRECT']}",
                size=16, bold=True, color=RED)
    add_textbox(slide, 0.7, 5.0, 12.0, 0.5,
                "Interpretation",
                size=14, bold=True, color=GREY_DARK)
    add_textbox(slide, 0.7, 5.5, 12.0, 1.5,
                "Point estimate favours memory by 5 pp, but with only 5 "
                "discordant pairs the sample size cannot confirm or refute. "
                "A bigger paired cohort is the obvious next step. The 100-"
                "patient aggregate on slide 3 mixes cohorts and is descriptive — "
                "this is the only controlled comparison.",
                size=12, color=GREY_DARK)


def slide_dashboard_overview(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Doctor dashboard — features overview",
               "Four features that make the run reviewable")
    console = REPO / "docs" / "final_presentation" / "doctor_console.png"
    if console.exists():
        add_image(slide, console, 0.4, 1.2, width=8.0)
    else:
        add_textbox(slide, 0.5, 1.5, 8.0, 1.0,
                    "[screenshot missing: docs/final_presentation/doctor_console.png]",
                    size=12, color=RED)
    items = [
        ("Agent workflow inline",
         "Every stage's output one click away, rendered as a doctor-readable narrative."),
        ("Similar past cases (Tier-4 recall)",
         "Top-K neighbours with their evaluator match type; one click switches view."),
        ("Treatment safety panel",
         "Drugs, interactions, contraindications + planner's assumptions & missing-data warnings."),
        ("Reviewer note + persistence",
         "Three-way verdict, free text, initials → data/gold/annotations/<uuid>.json."),
    ]
    left = 8.7
    top0 = 1.2
    for i, (head, body) in enumerate(items, start=1):
        top = top0 + (i - 1) * 1.25
        dia = 0.45
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(left), Inches(top),
                                         Inches(dia), Inches(dia))
        circle.fill.solid(); circle.fill.fore_color.rgb = AMBER
        circle.line.color.rgb = WHITE; circle.line.width = Pt(1.5)
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = str(i)
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
        add_textbox(slide, left + dia + 0.1, top - 0.05, 4.2, 0.45,
                    head, size=12, bold=True, color=NAVY)
        add_textbox(slide, left + dia + 0.1, top + 0.35, 4.2, 0.85,
                    body, size=10, color=GREY_DARK)
    add_textbox(slide, 0.4, 6.85, 12.5, 0.4,
                "URL-driven state: ?r=<set>&p=<uuid>&a=<agent> makes every view shareable and refresh-safe.",
                size=10, color=GREY_MED, align=PP_ALIGN.CENTER)


def slide_dashboard_treatment(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Treatment safety + reviewer flow",
               "What the planner did NOT know, surfaced; doctor verdict, persisted")
    shot = HERE / "dashboard_treatment.png"
    if shot.exists():
        add_image(slide, shot, 0.4, 1.2, width=8.5)
        text_left = 9.1
        text_width = 4.0
    else:
        add_textbox(slide, 0.4, 1.3, 8.5, 0.6,
                    "[treatment screenshot to be captured — drop dashboard_treatment.png next to build_pptx.py and rerun]",
                    size=11, color=RED)
        text_left = 0.4
        text_width = 12.5
    add_textbox(slide, text_left, 1.9 if not shot.exists() else 1.3,
                text_width, 0.5,
                "Surfaces what the planner did NOT know",
                size=14, bold=True, color=NAVY)
    add_textbox(slide, text_left, 2.4 if not shot.exists() else 1.8,
                text_width, 1.5,
                "• Drug dose assumptions (e.g. 'eGFR unknown — assumed normal "
                "for ACE-I dosing').\n"
                "• Missing-comorbidity warnings.\n"
                "• Interaction checks against current medication list.",
                size=11, color=GREY_DARK)
    add_textbox(slide, text_left, 4.0 if not shot.exists() else 3.6,
                text_width, 0.5,
                "Reviewer verdict persistence",
                size=14, bold=True, color=NAVY)
    add_textbox(slide, text_left, 4.5 if not shot.exists() else 4.1,
                text_width, 1.8,
                "• Three-way verdict (agree / uncertain / disagree).\n"
                "• Free text + reviewer initials.\n"
                "• Written to data/gold/annotations/<uuid>.json — the only "
                "write surface in the UI.\n"
                "• Unlocks clinician-agreement metrics beyond LLM-judge.",
                size=11, color=GREY_DARK)


def slide_literature(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Where CMADS sits against the literature",
               "Seven comparators · what they do, their result, and what they leave out")

    # Each item: (paper, what-they-do, result, missing-that-CMADS-adds, accent)
    items = [
        ("AMIE (Tu 2024)",
         "Conversational diagnostic AI; OSCE-style consultations.",
         "Beat PCPs on 28/32 specialist axes (149 cases · 20 PCPs).",
         "Closed-source backbone · no doctor UI · no inspectable memory.",
         BLUE),
        ("AgentClinic (Schmidgall 2024)",
         "Multimodal benchmark of LLM agents in simulated clinics.",
         "Claude-3.5: 62.1% on AgentClinic-MedQA · PCPs 54%±28.5.",
         "Benchmark of others, not a deployed pipeline · no treatment plan.",
         TEAL),
        ("MedAgents (Tang 2024, ACL)",
         "Multi-agent CoT for zero-shot medical reasoning.",
         "GPT-4: 86.7% avg on 9 MCQ benchmarks (MedQA 83.7%).",
         "MCQ-only · no patient data · no memory · no UI.",
         NAVY),
        ("MDTeamGPT (Chen 2025)",
         "Self-evolving MDT with CorrectKB + ChainKB.",
         "90.1% on MedQA · 83.9% on PubMedQA.",
         "MCQ-only · no end-to-end pipeline · no clinician interface.",
         AMBER),
        ("TAO (Kim 2025, ICML)",
         "Tiered Agentic Oversight for healthcare safety.",
         "+8.2% on 4/5 safety benchmarks · triage 40% → 60% w/ MDs.",
         "Safety wrapper, not a diagnostic pipeline · no case-based memory.",
         GREEN),
        ("ClinicalLab (Yan 2024)",
         "Coordinator + Chief Physician across 11 departments.",
         "Within ~5% of senior physicians on 1,500 real cases.",
         "Closed-source · no published memory subsystem · no doctor UI.",
         BLUE),
        ("ZODIAC (Zhou 2024)",
         "5-agent cardiology + clinician review.",
         "Cardiologist-level on 7/8 metrics (qualitative).",
         "Single specialty · no memory · no treatment plan · no open backbone.",
         TEAL),
    ]

    # 2-column grid: 4 left, 3 right
    col_w = 6.15
    row_h = 1.40
    rows = [(items[:4], 0.35), (items[4:], 6.85)]
    for col_items, left in rows:
        for j, (paper, what, result, missing, accent) in enumerate(col_items):
            top = 1.1 + j * row_h
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(left), Inches(top),
                                          Inches(col_w), Inches(row_h - 0.1))
            box.fill.solid(); box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = accent; box.line.width = Pt(1.25)
            # Paper title
            add_textbox(slide, left + 0.15, top + 0.05, col_w - 0.25, 0.36,
                        paper, size=12, bold=True, color=accent)
            # What they do
            add_textbox(slide, left + 0.15, top + 0.38, col_w - 0.25, 0.35,
                        what, size=9, color=GREY_DARK)
            # Result
            add_textbox(slide, left + 0.15, top + 0.70, col_w - 0.25, 0.35,
                        "Result: " + result, size=9, bold=True, color=GREY_DARK)
            # Missing
            add_textbox(slide, left + 0.15, top + 1.00, col_w - 0.25, 0.35,
                        "Missing: " + missing, size=9, color=RED)

    add_textbox(slide, 0.35, 6.85, 12.6, 0.4,
                "Comparison is on multiple axes: cohort + headline + what each system leaves out. CMADS fills the union of those gaps.",
                size=10, color=GREY_MED, align=PP_ALIGN.CENTER)


def slide_gaps(prs, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "Gaps I filled",
               "Each gap is small. Nobody closes all four in one system — that's the contribution.")
    quads = [
        ("1. End-to-end pipeline",
         "Synthea → differential → NICE plan → doctor review.",
         "vs MDAgents, ClinicalLab (benchmark accuracy only).",
         BLUE),
        ("2. Open-source backbone",
         "GPT-OSS-120B reasoning + Qwen3-32B judge, reproducible.",
         "vs MDAgents, ClinicalLab, ZODIAC (GPT-4 family).",
         TEAL),
        ("3. Inspectable 4-tier memory + controlled A/B",
         "Working / Episodic / Semantic / Case-based, with paired test.",
         "vs all five (memory implicit or absent).",
         AMBER),
        ("4. Doctor console + persisted verdicts",
         "Agent inspector, similar cases, treatment safety, annotation.",
         "vs all five (none ship a clinician annotation surface).",
         GREEN),
    ]
    card_w = 5.9
    card_h = 2.6
    positions = [(0.6, 1.3), (6.8, 1.3), (0.6, 4.1), (6.8, 4.1)]
    for (head, body, contrast, color), (left, top) in zip(quads, positions):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(top),
                                      Inches(card_w), Inches(card_h))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color; box.line.width = Pt(2)
        add_textbox(slide, left + 0.2, top + 0.15, card_w - 0.3, 0.55,
                    head, size=16, bold=True, color=color)
        add_textbox(slide, left + 0.2, top + 0.85, card_w - 0.3, 0.85,
                    body, size=12, color=GREY_DARK)
        add_textbox(slide, left + 0.2, top + 1.85, card_w - 0.3, 0.55,
                    contrast, size=11, color=GREY_MED)
    add_textbox(slide, 0.6, 6.85, 12.1, 0.45,
                "The composition is the contribution.",
                size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def main():
    metrics = compute_metrics.compute(REPO)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for fn in (slide_title, slide_what_is_cmads, slide_headline,
               slide_dashboard_overview, slide_literature, slide_gaps):
        fn(prs, metrics)

    prs.save(str(OUT))
    print(f"wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
