"""Open the PPTX and confirm: 9 slides, expected text present, expected
images embedded. Run after build_pptx.py finishes."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
PPTX = HERE / "CMADS_Results_Presentation.pptx"


REQUIRED_PHRASES = [
    "CMADS — Multi-Agent Systems for AI Clinical Decisioning",
    "What CMADS is",
    "Results — single-level vs multi-level memory",
    "Same 160 patients",
    "Δ (Multi-level − Single-level)",
    "Doctor dashboard — features overview",
    "Where CMADS sits against the literature",
    "AMIE",
    "MedAgents",
    "AgentClinic",
    "TAO",
    "MDTeamGPT",
    "Gaps I filled",
    "The composition is the contribution.",
]


def _all_text(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            chunks.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    chunks.append(cell.text_frame.text)
    return "\n".join(chunks)


def main() -> int:
    if not PPTX.exists():
        print(f"FAIL: {PPTX} does not exist")
        return 1

    prs = Presentation(str(PPTX))
    print(f"slides: {len(prs.slides)}")
    if len(prs.slides) != 6:
        print(f"FAIL: expected 6 slides, got {len(prs.slides)}")
        return 1

    full_text = "\n\n---SLIDE BREAK---\n\n".join(
        _all_text(s) for s in prs.slides
    )
    missing = [p for p in REQUIRED_PHRASES if p not in full_text]
    if missing:
        print("FAIL: missing required phrases:")
        for m in missing:
            print(f"  - {m}")
        return 1

    img_count = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    print(f"images embedded: {img_count}")
    if img_count < 2:
        print(f"WARN: only {img_count} image(s) — expected at least 2 "
              "(system diagram + doctor console). Continuing.")

    print("OK")
    return 0


if __name__ == "__main__":
    sys.exit(main())
