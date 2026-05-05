"""Build the dedicated Multi-Level Memory deck (CMADS).

Produces CMADS_Multi_Level_Memory.pptx — a deeper deck than the supervisor
progress slide, walking through:

    1. Title
    2. The problem the old shared state hid
    3. Old vs new memory side-by-side
    4. The 4-tier architecture in detail
    5. Tier 1 (working) + Tier 2 (episodic)
    6. Tier 3 (semantic) + Tier 4 (procedural)
    7. How memory flows through the pipeline (re-uses agents_block_diagram)
    8. A/B results — Memory ON vs OFF on Batch 4
    9. Implementation summary + next steps

Diagrams used:
    - memory_old_vs_new.png  (built from memory_old_vs_new.svg)
    - memory_architecture.png (built from memory_architecture.svg)
    - agents_block_diagram.png (re-used from the progress deck)

Run:
    python3 docs/memory_presentation/build_memory_pptx.py
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
PROGRESS = HERE.parent / "progress_presentation"
ARCH_SVG = HERE / "memory_architecture.svg"
ARCH_PNG = HERE / "memory_architecture.png"
COMPARE_SVG = HERE / "memory_old_vs_new.svg"
COMPARE_PNG = HERE / "memory_old_vs_new.png"
AGENTS_PNG = PROGRESS / "agents_block_diagram.png"
OUT_PPTX = HERE / "CMADS_Multi_Level_Memory.pptx"

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"


# ── Colours (matches the progress deck) ────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x16, 0xA0, 0x85)
AMBER = RGBColor(0xD3, 0x54, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
GREY_DARK = RGBColor(0x34, 0x49, 0x5E)
GREY_MED = RGBColor(0x7F, 0x8C, 0x8D)
GREY_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)
WARM = RGBColor(0xFE, 0xF9, 0xE7)


# ── Tier accent colours ────────────────────────────────────
T1 = AMBER
T2 = BLUE
T3 = TEAL
T4 = NAVY


# ── PNG rendering helper ───────────────────────────────────
def ensure_png(svg: Path, png: Path, w: int = 1600, h: int = 760) -> Path:
    """Render an SVG to PNG via headless Chrome (skips if PNG already exists)."""
    if png.exists():
        return png
    subprocess.run(
        [
            CHROME, "--headless", "--disable-gpu",
            "--force-device-scale-factor=2",
            f"--screenshot={png}",
            f"--window-size={w},{h}",
            "--hide-scrollbars",
            "--default-background-color=FFFFFFFF",
            svg.as_uri(),
        ],
        check=True, capture_output=True,
    )
    return png


# ── Generic helpers ────────────────────────────────────────

def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_header_bar(slide, title: str, subtitle: str | None = None,
                   accent: RGBColor = NAVY):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.7)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.name = "Calibri"; r.font.size = Pt(26)
    r.font.bold = True; r.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = "Calibri"; r2.font.size = Pt(12)
        r2.font.color.rgb = GREY_LIGHT
        r2.font.italic = True


def add_textbox(slide, left, top, width, height, text, *,
                size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                italic=False, name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = name; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items,
                size=13, color=GREY_DARK, gap=Pt(6), bullet_char="•  "):
    """items: list[str | (text, bold) | (text, bold, color)]."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            if len(item) == 3:
                text, bold, c = item
            else:
                text, bold = item
                c = color
        else:
            text, bold, c = item, False, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = gap
        r = p.add_run(); r.text = bullet_char + text
        r.font.name = "Calibri"; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = c
    return tb


def add_footer(slide, page: str):
    add_textbox(slide, Inches(0.3), Inches(7.1), Inches(8), Inches(0.3),
                "CMADS · Multi-Level Memory · Bachelor Thesis · Abdelrahman",
                size=10, color=GREY_MED, italic=True)
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
                page, size=10, color=GREY_MED, align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=NAVY, card_color=WHITE):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid(); card.fill.fore_color.rgb = card_color
    card.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tb = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.name = "Calibri"; r.font.size = Pt(28); r.font.bold = True
    r.font.color.rgb = value_color
    tb2 = slide.shapes.add_textbox(left, top + Inches(0.85), width, Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.name = "Calibri"; r2.font.size = Pt(11); r2.font.color.rgb = GREY_DARK


def add_tier_card(slide, left, top, width, height, tier_no, name, scope,
                  body_lines, code_lines, accent: RGBColor):
    """A vertically-stacked card describing one memory tier."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent; card.line.width = Pt(1.4)
    card.shadow.inherit = False

    # Accent header strip
    head_h = Inches(0.55)
    head = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, head_h
    )
    head.fill.solid(); head.fill.fore_color.rgb = accent
    head.line.fill.background()

    add_textbox(slide, left + Inches(0.2), top + Inches(0.08),
                width - Inches(0.4), Inches(0.35),
                f"TIER {tier_no} · {name}",
                size=14, bold=True, color=WHITE)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.32),
                width - Inches(0.4), Inches(0.3),
                scope, size=10, italic=True, color=GREY_LIGHT)

    # Body
    body_top = top + head_h + Inches(0.18)
    add_bullets(slide, left + Inches(0.25), body_top,
                width - Inches(0.5), Inches(2.5),
                body_lines, size=12, color=GREY_DARK, gap=Pt(4),
                bullet_char="·  ")

    # Code box at the bottom
    code_top = top + height - Inches(1.05)
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(0.2), code_top,
        width - Inches(0.4), Inches(0.95)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    code_box.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
    code_box.line.width = Pt(0.5)
    code_tb = slide.shapes.add_textbox(
        left + Inches(0.3), code_top + Inches(0.08),
        width - Inches(0.6), Inches(0.85)
    )
    tf = code_tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run(); r.text = line
        r.font.name = "Menlo"; r.font.size = Pt(10.5)
        r.font.color.rgb = GREY_DARK


# ── Build the deck ─────────────────────────────────────────
def build():
    ensure_png(ARCH_SVG, ARCH_PNG, w=1600, h=760)
    ensure_png(COMPARE_SVG, COMPARE_PNG, w=1600, h=760)
    if not AGENTS_PNG.exists():
        # Render the agents block diagram if it isn't there yet.
        ensure_png(PROGRESS / "agents_block_diagram.svg",
                   AGENTS_PNG, w=1600, h=700)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1 · Title ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, PAPER)

    band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), Inches(7.5)
    )
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.4), Inches(11.5), Inches(0.5),
                "BACHELOR THESIS · DEEP DIVE", size=14,
                bold=True, color=AMBER)

    add_textbox(s, Inches(1), Inches(1.95), Inches(11.5), Inches(1.4),
                "Multi-Level Memory", size=58, bold=True, color=NAVY)

    add_textbox(s, Inches(1), Inches(3.3), Inches(11.5), Inches(0.9),
                "Architecture, Implementation, Impact",
                size=28, bold=True, color=BLUE)

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.5), Inches(0.7),
                "How CMADS agents share session context — supervisor request, "
                "May 2026.",
                size=15, italic=True, color=GREY_DARK)

    line = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.1),
        Inches(3), Inches(0.04)
    )
    line.fill.solid(); line.fill.fore_color.rgb = AMBER
    line.line.fill.background()

    add_textbox(s, Inches(1), Inches(5.3), Inches(11.5), Inches(0.4),
                "Presented by  Abdelrahman", size=15,
                bold=True, color=GREY_DARK)
    add_textbox(s, Inches(1), Inches(5.75), Inches(11.5), Inches(0.4),
                "Supervisor follow-up · CMADS Project",
                size=13, color=GREY_MED)

    add_textbox(
        s, Inches(1), Inches(6.55), Inches(11.5), Inches(0.3),
        "Agenda  ·  the gap in the old design  ·  4-tier architecture  ·  "
        "tier-by-tier walk-through  ·  pipeline integration  ·  A/B results",
        size=11, color=GREY_MED, italic=True,
    )

    add_speaker_notes(s, """
This is a deeper follow-up to last week's progress meeting where you asked
for a multi-level memory system. The next eight slides cover the gap in the
old design, the four-tier architecture I built, how each tier works, how it
plugs into the pipeline, and what the before-and-after experiment showed.
""")

    # ── Slide 2 · The problem ─────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "The gap in the old shared state",
        "Why agents needed more than each other's final JSON outputs"
    )

    # Left: scenario
    add_textbox(s, Inches(0.4), Inches(1.05), Inches(7), Inches(0.4),
                "A concrete scenario", size=15, bold=True, color=AMBER)

    scenario = (
        "The Diagnostic Reasoning agent runs an adaptive critique loop:\n"
        "  · round 1 confidence = 60 (rejected: anchoring on lab finding)\n"
        "  · round 2 confidence = 72 (added missed dx, recalibrated)\n"
        "  · round 3 confidence = 78 → stops, emits final differential\n"
        "\nThe Reviewer's job is to challenge that final differential. "
        "But the old state exposes ONLY the final JSON. The reasoning trail "
        "is gone — the Reviewer can't see WHICH hypotheses were considered "
        "and rejected, what the confidence trajectory was, or where the "
        "Diagnostic agent might have anchored too quickly."
    )
    tb = s.shapes.add_textbox(Inches(0.4), Inches(1.45), Inches(7),
                              Inches(3.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(scenario.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        r = p.add_run(); r.text = line
        r.font.name = "Calibri"; r.font.size = Pt(13)
        r.font.color.rgb = GREY_DARK

    # Bottom-left: limits
    add_textbox(s, Inches(0.4), Inches(5.1), Inches(7), Inches(0.4),
                "Concretely, the old design hides:",
                size=14, bold=True, color=AMBER)
    add_bullets(
        s, Inches(0.4), Inches(5.5), Inches(7), Inches(1.6),
        [
            "the Diagnostic agent's reasoning rounds + confidence trajectory",
            "the hypotheses considered and rejected within a run",
            "any signal from one patient run into the next",
            "aggregate priors over diseases observed in past runs",
        ],
        size=12, color=GREY_DARK, gap=Pt(4), bullet_char="✗  ",
    )

    # Right: panel
    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.7), Inches(1.05), Inches(5.3), Inches(5.55)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    panel.line.color.rgb = RED; panel.line.width = Pt(1)
    panel.shadow.inherit = False

    add_textbox(s, Inches(7.9), Inches(1.2), Inches(5), Inches(0.4),
                "The old PipelineState — at a glance",
                size=14, bold=True, color=RED)
    add_bullets(
        s, Inches(7.9), Inches(1.7), Inches(5), Inches(4.7),
        [
            ("patient_context — set once, read by all", True),
            ("agent_outputs — per-agent JSON slot", True),
            ("execution_trace — status / timing only", True),
            ("conflicts — declared, NEVER written", False, GREY_MED),
            ("scratchpad — declared, NEVER written", False, GREY_MED),
            ("", False),
            ("3 truly active channels.", True, RED),
            ("0 cross-session memory.", True, RED),
            ("0 typed reasoning timeline.", True, RED),
        ],
        size=12, color=GREY_DARK, gap=Pt(5), bullet_char="·  ",
    )

    add_footer(s, "2 / 9")
    add_speaker_notes(s, """
Concretely: the Diagnostic agent's adaptive loop runs three rounds with
shifting confidence — but the Reviewer only sees the final differential JSON.
The reasoning chain is gone. And once a patient run finishes, the system
has no way to remember anything about it for the next patient. The old
PipelineState had five channels but two were declared and never written —
they were the placeholder shape of a memory system without the substance.
""")

    # ── Slide 3 · Old vs New (diagram) ─────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "Old memory vs. new memory",
        "Same LangGraph state — but every channel now load-bearing, plus three "
        "new layers above and beyond"
    )

    # Place the comparison PNG centrally
    img_w = Inches(12.6)
    img_h = img_w * (760 / 1600)  # ≈ 5.99 in
    s.shapes.add_picture(str(COMPARE_PNG),
                         Inches((13.333 - 12.6) / 2),
                         Inches(1.05),
                         width=img_w, height=img_h)

    add_footer(s, "3 / 9")
    add_speaker_notes(s, """
Left side, the old design: five channels in PipelineState, three actually
used. Right side, the new design: every channel is load-bearing —
scratchpad becomes Tier 1 working memory, two new channels handle Tier 2
episodic, and two storage layers live outside the state for Tier 3 and
Tier 4. The conflicts channel stays reserved for forward compatibility but
is still unused — I didn't want to remove a public state field while it's
in the shared interface.
""")

    # ── Slide 4 · 4-tier architecture ─────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "Four-tier memory architecture",
        "CoALA-inspired · accessed through one MemoryManager facade · "
        "no LLM coupling — fast, deterministic recall"
    )

    img_w = Inches(12.6)
    img_h = img_w * (760 / 1600)
    s.shapes.add_picture(str(ARCH_PNG),
                         Inches((13.333 - 12.6) / 2),
                         Inches(1.05),
                         width=img_w, height=img_h)

    add_footer(s, "4 / 9")
    add_speaker_notes(s, """
The four tiers, scoped from narrowest to broadest. Working memory lives for
one agent invocation. Episodic memory lives for one patient run. Semantic
memory lives forever, on disk. Procedural memory is the long-term knowledge
store — the NICE clinical guidelines we already had in Qdrant. Every agent
talks to all four through a single MemoryManager facade, so the agent
code stays clean and the recall semantics are uniform across tiers.
""")

    # ── Slide 5 · Tier 1 + Tier 2 deep dive ──────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "Tier 1 · Working   /   Tier 2 · Episodic",
        "Inside the run — per-agent scratch and the reasoning timeline"
    )

    # Tier 1 card (left)
    add_tier_card(
        s, Inches(0.4), Inches(1.05), Inches(6.25), Inches(5.7),
        tier_no=1, name="WORKING MEMORY",
        scope="per-agent · invocation-local · in-state",
        body_lines=[
            "Lives in state.scratchpad[agent_id], "
            "previously a declared-but-unused channel.",
            "Used heavily by the Diagnostic agent's adaptive loop to track "
            "confidence trajectory + critique trail across rounds.",
            "Lets the agent's own follow-up calls reason about what already "
            "happened, instead of reconcatenating free text.",
            "Exposed via mm.working.put() / get() / append_to() / snapshot().",
            "Cleared at the start of each invocation — no leakage across "
            "agent calls.",
        ],
        code_lines=[
            "mm = MemoryManager.from_state(state, self.agent_id)",
            "mm.working.put('confidence_trajectory', [60, 72, 78])",
            "mm.working.append_to('critique_trail',",
            "    {'round': 2, 'confidence': 72, 'adequate': False})",
        ],
        accent=T1,
    )

    # Tier 2 card (right)
    add_tier_card(
        s, Inches(6.85), Inches(1.05), Inches(6.1), Inches(5.7),
        tier_no=2, name="EPISODIC MEMORY",
        scope="current run · in-state, persisted to disk",
        body_lines=[
            "New state channel session_memory — append-reducer, list of typed "
            "SessionEvent dicts (Pydantic-validated).",
            "Event types: agent_start / agent_complete / critique / "
            "confidence_check / decision / hypothesis / evidence_link.",
            "Reviewer reads diagnostic_reasoning's critique events to challenge "
            "the path to the diagnosis, not just the final JSON.",
            "Refiner reads the full session timeline + cross-session priors "
            "before producing the final differential.",
            "Saved on disk per-run as session_memory.json next to the existing "
            "execution_trace.json.",
        ],
        code_lines=[
            "self._pending_memory_events.append(",
            "  EpisodicMemory.write(",
            "    event_type='critique', agent_id=self.agent_id,",
            "    summary=f'Round {n} conf={c} adequate={a}',",
            "    payload={'round': n, 'confidence': c}))",
        ],
        accent=T2,
    )

    add_footer(s, "5 / 9")
    add_speaker_notes(s, """
Tier 1 turns the previously-unused scratchpad channel into real working
memory for multi-call agents. Tier 2 is genuinely new: an append-reducer
list of typed SessionEvents that lets downstream agents read the WHY behind
each decision, not just the WHAT. The Reviewer is the biggest beneficiary —
it can now cite specific critique rounds when challenging a diagnosis.
""")

    # ── Slide 6 · Tier 3 + Tier 4 deep dive ──────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "Tier 3 · Semantic   /   Tier 4 · Procedural",
        "Beyond the run — cross-session priors and long-term knowledge"
    )

    # Tier 3 card
    add_tier_card(
        s, Inches(0.4), Inches(1.05), Inches(6.25), Inches(5.7),
        tier_no=3, name="SEMANTIC MEMORY",
        scope="cross-session · on disk · per-disease aggregate",
        body_lines=[
            "Lives in data/gold/memory/semantic_memory.json — a small JSON file "
            "keyed by disease name (no embeddings at this scale).",
            "Stores per-disease counts (DIRECT / INDIRECT / MISS), rank-1 "
            "frequency, average primary confidence, observed evidence patterns.",
            "Updated by the new Stage-7 memory_consolidation_node after every "
            "patient run — running mean for the average, dedup for the "
            "evidence patterns, capped to a fixed max.",
            "Recall is case-insensitive, idempotent, atomic write via tmp-file "
            "+ replace; resilient to corrupt files (returns None, doesn't crash).",
            "Read by the Refiner and Treatment agents as Bayesian-style priors "
            "for the candidate diseases.",
        ],
        code_lines=[
            "store = SemanticMemory(cfg.SEMANTIC_MEMORY_PATH)",
            "insight = store.recall('Hypertension')",
            "# {runs:14, direct:9, rank1:7, avg_conf:81.2}",
            "store.consolidate('Hypertension', 'DIRECT', rank=1, ...)",
        ],
        accent=T3,
    )

    # Tier 4 card
    add_tier_card(
        s, Inches(6.85), Inches(1.05), Inches(6.1), Inches(5.7),
        tier_no=4, name="PROCEDURAL MEMORY",
        scope="long-term · static · NICE guidelines via Qdrant",
        body_lines=[
            "Wraps the existing NICE-guidelines vector store in Qdrant — no "
            "new infra, just a uniform recall(query) API symmetric to the "
            "other tiers.",
            "Same code path the Treatment agent already used; now also "
            "reachable from Diagnostic and Reviewer if they want guideline "
            "context for differential validation.",
            "Graceful degradation: if Qdrant is unreachable (no QDRANT_URL, "
            "network down) the wrapper returns [] instead of crashing the run.",
            "Embedding model unchanged: BioLORD-2023 over the existing "
            "nice_guidelines collection — semantic search by disease name.",
            "Static knowledge — never updated by the pipeline. Curated by hand.",
        ],
        code_lines=[
            "guideline = mm.procedural.lookup_disease('Hypertension')",
            "# {nice_guideline:'NG136',",
            "#  recommended_drugs:['ACE-i','CCB','...']}",
            "results = mm.procedural.recall('proteinuria', top_k=3)",
        ],
        accent=T4,
    )

    add_footer(s, "6 / 9")
    add_speaker_notes(s, """
Tier 3 is where one run's outcome flows into the next. After each patient,
the consolidation node writes a small aggregate record per matched disease:
how often it landed at rank 1, what evidence patterns the agent attached
to it, the running average primary confidence. The Refiner reads this on
the next run as a prior. Tier 4 is the existing Qdrant store of NICE
guidelines — wrapped here so all four tiers feel the same to the agent
code, with graceful degradation if the vector store is unreachable.
""")

    # ── Slide 7 · How it flows in the pipeline ────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "How memory flows through the pipeline",
        "Same 7 agents · 1 new maintenance node · per-agent reads/writes"
    )

    # Top half: agents block diagram
    img_w = Inches(12.4)
    img_h = img_w * (700 / 1600)  # original aspect ratio
    s.shapes.add_picture(str(AGENTS_PNG),
                         Inches((13.333 - 12.4) / 2),
                         Inches(1.05),
                         width=img_w, height=img_h)

    # Bottom half: per-agent matrix
    matrix_top = Inches(1.1) + img_h + Inches(0.1)
    matrix_left = Inches(0.4)
    matrix_w = Inches(12.55)
    matrix_h = Inches(7.5) - matrix_top - Inches(0.5)

    rows = [
        ("Agent",                "Reads",                    "Writes"),
        ("EHR Analyst",          "—",                        "T2 lifecycle"),
        ("Lab Interpreter",      "—",                        "T2 lifecycle"),
        ("Diagnostic Reasoning", "T2 lifecycle (own)",       "T1 trajectory · T2 critique·conf·decision"),
        ("Clinical Reviewer",    "T2 critique trail (Diag)", "T2 lifecycle"),
        ("Refiner (final dx)",   "T2 timeline · T3 priors",  "T2 decision"),
        ("Treatment",            "T3 priors · T4 guidelines","T2 decision"),
        ("Memory Consolidation", "T2 timeline · final dx",   "T3 disease record (NEW)"),
    ]
    table = s.shapes.add_table(
        len(rows), 3,
        matrix_left, matrix_top, matrix_w, matrix_h,
    ).table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(4.8)
    table.columns[2].width = Inches(5.15)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(11)
                    r.font.bold = (i == 0)
                    if i == 0:
                        r.font.color.rgb = WHITE
                    elif i == len(rows) - 1:
                        r.font.color.rgb = WHITE
                    else:
                        r.font.color.rgb = GREY_DARK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif i == len(rows) - 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = AMBER
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(
                    0xF4, 0xF6, 0xF7
                )

    add_footer(s, "7 / 9")
    add_speaker_notes(s, """
The pipeline shape is unchanged at the top — same seven agents. What's
new is the read/write matrix below: every agent now writes lifecycle events
to Tier 2; the Diagnostic agent also writes the per-round critique and
confidence checks; the Reviewer reads that trail; the Refiner reads the
whole timeline plus Tier-3 priors for the candidate diseases. After
Treatment, a new Stage-7 node consolidates the run into Tier 3 so the
next run starts smarter.
""")

    # ── Slide 8 · A/B Results ─────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "A/B results — Memory ON vs OFF",
        "20 patients · Batch 4 · same model (GPT-OSS 120B) · "
        "same judge (Qwen3 32B) · only MEMORY_ENABLED toggled"
    )

    # KPI strip
    kpi_top = Inches(1.15)
    kpi_w = Inches(2.85); kpi_h = Inches(1.45)
    add_kpi_card(s, Inches(0.4),  kpi_top, kpi_w, kpi_h, "+5 pp",
                 "DIRECT match (40 % → 45 %)", value_color=GREEN)
    add_kpi_card(s, Inches(3.45), kpi_top, kpi_w, kpi_h, "+10 pp",
                 "Found rate (80 % → 90 %)", value_color=GREEN)
    add_kpi_card(s, Inches(6.5),  kpi_top, kpi_w, kpi_h, "≈ 0",
                 "Wall-clock cost (113 → 114 s)", value_color=AMBER)
    add_kpi_card(s, Inches(9.55), kpi_top, kpi_w, kpi_h, "−2 pp",
                 "Rank-1 within found (19 % → 17 %)", value_color=GREY_DARK)

    # Detailed table on the left
    rows = [
        ("Metric", "Memory OFF", "Memory ON", "Δ"),
        ("DIRECT match",        "8/20 · 40 %", "9/20 · 45 %", "+5 pp"),
        ("Found rate (D + I)",  "80 %",        "90 %",        "+10 pp"),
        ("Rank-1 when found",   "19 %",        "17 %",        "−2 pp"),
        ("Avg time / patient",  "113 s",       "114 s",       "≈ 0"),
        ("Patients evaluated",  "20",          "20",          "—"),
    ]
    tbl_left = Inches(0.4); tbl_top = Inches(2.85)
    tbl_w = Inches(7.4); tbl_h = Inches(3.0)
    table = s.shapes.add_table(
        len(rows), 4, tbl_left, tbl_top, tbl_w, tbl_h
    ).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(1.7)
    table.columns[2].width = Inches(1.7)
    table.columns[3].width = Inches(1.6)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(12)
                    r.font.bold = (i == 0)
                    if i == 0:
                        r.font.color.rgb = WHITE
                    elif j == 3 and ("+" in val):
                        r.font.color.rgb = GREEN
                        r.font.bold = True
                    elif j == 3 and "−" in val:
                        r.font.color.rgb = GREY_DARK
                    else:
                        r.font.color.rgb = GREY_DARK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(
                    0xF4, 0xF6, 0xF7
                )

    # Findings panel on the right
    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.0), Inches(2.85), Inches(5.0), Inches(3.95)
    )
    panel.fill.solid(); panel.fill.fore_color.rgb = WARM
    panel.line.color.rgb = RGBColor(0xD4, 0xAC, 0x0D); panel.line.width = Pt(1)
    panel.shadow.inherit = False

    add_textbox(s, Inches(8.2), Inches(2.97), Inches(4.6), Inches(0.45),
                "What the deltas mean", size=14, bold=True,
                color=RGBColor(0x7D, 0x66, 0x08))
    add_bullets(
        s, Inches(8.2), Inches(3.45), Inches(4.6), Inches(3.3),
        [
            ("Memory wiring helps detection more than ranking.", True),
            ("Found-rate gain comes from the Reviewer reading "
             "Diagnostic's critique trail and challenging earlier.", False),
            ("Rank-1-within-found dropped 2 pp — within noise on "
             "20 patients (one borderline case flips this).", False),
            ("No measurable time penalty (113 → 114 s) — recall is "
             "rule-based and deterministic, not LLM-driven.", False),
            ("Same model, same judge, same prompts — only "
             "MEMORY_ENABLED toggled. Clean A/B.", True),
        ],
        size=11.5, color=GREY_DARK, gap=Pt(5), bullet_char="·  ",
    )

    # Caveats strip
    add_textbox(
        s, Inches(0.4), Inches(6.0), Inches(12.6), Inches(0.45),
        "Caveats   ·   n=20 (small)   ·   single batch   ·   "
        "single LLM (GPT-OSS 120B)   ·   no model interaction effects "
        "tested yet   ·   semantic memory empty at start of run "
        "(no priors yet to recall)",
        size=10.5, italic=True, color=GREY_MED,
    )

    add_footer(s, "8 / 9")
    add_speaker_notes(s, """
The headline: five percentage points more DIRECT matches and ten more
percentage points found-rate, at no measurable time cost. Rank-1
within-found stays roughly flat — that one slipped two points but on
twenty patients that is one borderline case flipping. The way I read it,
memory wiring is helping the system DETECT more correct diagnoses,
because the Reviewer can now challenge the path more sharply. Caveats:
twenty patients is small; the semantic store starts empty so its priors
only kick in for runs after the first; same model and judge throughout.
The next step is to validate at full 270-patient cohort scale.
""")

    # ── Slide 9 · Summary + next steps ────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header_bar(
        s, "Summary  ·  Next steps",
        "Where the multi-level memory feature stands today"
    )

    # Big nav cards — what was delivered
    card_top = Inches(1.2); card_w = Inches(4.0); card_h = Inches(2.5)
    titles = ["Architecture", "Implementation", "Validation"]
    contents = [
        [
            "4 tiers (CoALA-inspired)",
            "1 MemoryManager facade",
            "Memory_consolidation node (Stage 7)",
            "Master flag MEMORY_ENABLED",
        ],
        [
            "src/memory/ package — 7 modules",
            "Wired into 4 agents (Diag/Rev/Ref/Tx)",
            "24 unit tests · all pass",
            "Backward-compatible state shape",
        ],
        [
            "20-patient A/B on Batch 4",
            "+5 pp DIRECT, +10 pp Found",
            "no time cost (113 → 114 s)",
            "memory_old_vs_new + memory_architecture diagrams",
        ],
    ]
    accents = [BLUE, AMBER, GREEN]
    for i, (title, items, accent) in enumerate(zip(titles, contents, accents)):
        left = Inches(0.4 + i * 4.3)
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_w, card_h
        )
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = accent; card.line.width = Pt(1.4)
        card.shadow.inherit = False
        # Header strip
        head = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, card_top, card_w, Inches(0.5)
        )
        head.fill.solid(); head.fill.fore_color.rgb = accent
        head.line.fill.background()
        add_textbox(s, left + Inches(0.2), card_top + Inches(0.08),
                    card_w - Inches(0.4), Inches(0.35),
                    title, size=14, bold=True, color=WHITE)
        add_bullets(
            s, left + Inches(0.25), card_top + Inches(0.65),
            card_w - Inches(0.5), card_h - Inches(0.7),
            items, size=12, color=GREY_DARK, gap=Pt(5),
            bullet_char="·  ",
        )

    # Next steps
    add_textbox(s, Inches(0.4), Inches(4.0), Inches(12.6), Inches(0.4),
                "Next steps", size=15, bold=True, color=NAVY)
    add_bullets(
        s, Inches(0.4), Inches(4.5), Inches(12.6), Inches(2.4),
        [
            ("Validate at full 270-patient cohort scale and quantify "
             "per-disease deltas (does the wiring help the weakest "
             "categories more — CKD staging, plain hypertension?)",
             False),
            ("Warm-start the semantic store from the existing 160-patient "
             "results so priors kick in from patient 1 of any new run.",
             False),
            ("Test memory wiring × LLM choice — does the gain hold for "
             "Med42 70B and Claude Sonnet, or is it model-specific?",
             False),
            ("Add memory recall events to the trace itself, so we can "
             "audit which prior fired in each agent decision.", False),
            ("Optional: relax DIRECT criterion alongside memory — "
             "the two effects should compound on the weakest diseases.",
             True),
        ],
        size=12, color=GREY_DARK, gap=Pt(7), bullet_char="✓  ",
    )

    # Closing strip
    add_textbox(s, Inches(0.4), Inches(6.7), Inches(12.6), Inches(0.45),
                "Thank you — questions?", size=20, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(s, "9 / 9")
    add_speaker_notes(s, """
Summary card on the left: four tiers, one facade, one new stage. Middle
card: a clean Python package, wired into the four agents that benefit,
twenty-four unit tests covering the contracts. Right card: validated on
a clean A/B on Batch 4 with the same model and judge. Next is to scale
this to the full cohort, warm-start the semantic store, and check that
the gain isn't model-specific. Happy to take questions.
""")

    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    build()
